"""V9 ("flag, never suppress") must reach the rendered deck.

Before this, V3-V6 notes reached the slide footer (see
tests/test_assemble_profile_auto_footnotes.py) but V9 -- the KG
negative-equity flag the real FY2024 Seidensticker Konzernbilanz actually
discloses -- had no path anywhere into render output: not a footnote (a
committed test explicitly forbids that conflation), not a block, nothing.
These tests render the real Ankura master end to end and read the produced
.pptx back, rather than only checking the intermediate profile JSON, so a
regression that silently drops the advisory line again cannot pass quietly.
"""
from __future__ import annotations

import json
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

from render.assemble_profile import assemble_profile
from render.renderer import SLIDE_INDEX, render_profile

E2E_PROFILE = json.loads((Path(__file__).parent / "fixtures" / "e2e_output" / "seidensticker_e2e.json").read_text(encoding="utf-8"))["profile"]


@pytest.fixture
def output_dir():
    # Not pytest's own tmp_path: its base-temp discovery enumerates a shared
    # C:\Users\<user>\AppData\Local\Temp\pytest-of-<user> directory that this
    # environment denies access to (reproduced twice, PermissionError
    # [WinError 5] on os.scandir -- unrelated to anything under test here).
    # tempfile.TemporaryDirectory() creates its own uniquely-named directory
    # without that enumeration step and works fine in the same environment.
    with tempfile.TemporaryDirectory() as directory:
        yield Path(directory)


def _footer_text(pptx_path: Path) -> str:
    presentation = Presentation(pptx_path)
    slide = presentation.slides[SLIDE_INDEX]
    footer = next(shape for shape in slide.shapes if shape.has_text_frame and shape.top / 914400 > 6.9)
    return footer.text_frame.text


def test_v9_reaches_the_footer_as_a_distinct_advisory_line(output_dir: Path) -> None:
    normalised = {"entity": E2E_PROFILE["entity"], "rows": E2E_PROFILE["rows"]}
    profile = assemble_profile(normalised)
    output = output_dir / "advisory.pptx"

    render_profile(profile, output)

    footer_text = _footer_text(output)
    assert "⚠ Advisory:" in footer_text
    assert "Negative equity position disclosed" in footer_text
    # Never conflated with the V3-V6 note mechanism's own line.
    note_line, advisory_line = footer_text.split("\n")[0], footer_text.split("\n")[-1]
    assert "Negative equity" not in note_line
    assert advisory_line != note_line

    companion = json.loads(output.with_suffix(".json").read_text(encoding="utf-8"))
    assert any("Negative equity" in line for line in companion["render"]["advisory"])

    notes_text = presentation_notes(output)
    assert "BS-P-NEGEQ" in notes_text


def presentation_notes(pptx_path: Path) -> str:
    presentation = Presentation(pptx_path)
    slide = presentation.slides[SLIDE_INDEX]
    return slide.notes_slide.notes_text_frame.text


def test_no_advisory_line_when_v9_does_not_fire(output_dir: Path) -> None:
    rows = [row for row in E2E_PROFILE["rows"] if row["std_id"] != "BS-P-NEGEQ"]
    normalised = {"entity": E2E_PROFILE["entity"], "rows": rows}
    profile = assemble_profile(normalised)
    output = output_dir / "no_advisory.pptx"

    render_profile(profile, output)

    footer_text = _footer_text(output)
    assert "Advisory" not in footer_text

    companion = json.loads(output.with_suffix(".json").read_text(encoding="utf-8"))
    assert companion["render"]["advisory"] == []
