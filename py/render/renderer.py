"""Render a four-box company profile into the Ankura reference master.

The reference deck's financial annotations are floating shapes, rather than
chart data labels.  This module deliberately rebuilds both the native chart and
every displayed annotation from the same series, so a changed bar can never
retain an old printed value.
"""
from __future__ import annotations

import json
import shutil
from collections.abc import Iterable, Mapping, Sequence
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from render.assemble_profile import flagged_items
from validate.validator import validate_normalised, validate_v11

MASTER_FILENAME = "ankura_master_reference.pptx"
SLIDE_INDEX = 5  # Zero-based: slide 6, the designated At a glance slide.


class RenderError(ValueError):
    """Raised when a profile cannot be rendered without making up content."""


@dataclass(frozen=True)
class Slot:
    name: str
    x: float
    y: float
    width: float = 5.97
    height: float = 2.67


SLOTS: dict[str, Slot] = {
    "top_left": Slot("top_left", 0.42, 1.27),
    "top_right": Slot("top_right", 6.91, 1.27),
    "bottom_left": Slot("bottom_left", 0.46, 4.12),
    "bottom_right": Slot("bottom_right", 6.90, 4.12),
}


def render_profile(
    profile: Mapping[str, Any],
    output_pptx: str | Path,
    *,
    template: str | Path | None = None,
    assignments: Mapping[str, str] | None = None,
    human_notes: Mapping[str, Sequence[str]] | Sequence[str] | None = None,
) -> Path:
    """Render ``profile`` and its audit companion next to ``output_pptx``.

    ``assignments`` is intentionally required either as an argument or as
    ``profile['slot_assignments']`` / ``profile['canonical_layout']``.  The
    renderer never substitutes a different block for a missing selection.
    """
    resolved_template = _resolve_template(template)
    output = Path(output_pptx).resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    selected = dict(assignments or profile.get("slot_assignments") or profile.get("canonical_layout") or {})
    _validate_assignments(profile, selected)

    # Start with a byte-for-byte copy so layouts, master, notes, and unknown
    # OLE parts (including think-cell) survive. python-pptx preserves these
    # unmodelled package parts on save.
    shutil.copyfile(resolved_template, output)
    presentation = Presentation(output)
    if len(presentation.slides) <= SLIDE_INDEX:
        raise RenderError("Ankura master does not contain slide 6.")
    slide = presentation.slides[SLIDE_INDEX]
    if slide.slide_layout.name != "1_Blank with Title":
        raise RenderError(
            f"Slide 6 must use layout '1_Blank with Title', got {slide.slide_layout.name!r}."
        )

    blocks = {str(block["id"]): block for block in profile["blocks"]}
    _clear_content_regions(slide)
    _set_heading(slide, profile["entity"], blocks, selected)

    audit_entries: list[dict[str, Any]] = []
    rendered_footnotes: list[str] = []
    sources: list[str] = []
    for slot_name, block_id in selected.items():
        block = blocks[block_id]
        slot = SLOTS[slot_name]
        _render_block(slide, slot, block)
        audit_entries.extend(_audit_entries(slot_name, block))
        rendered_footnotes.extend(str(note) for note in block.get("footnotes_auto", []))
        rendered_footnotes.extend(_flag_notes(block))
        sources.extend(_source_lines(block))

    rendered_footnotes.extend(_human_notes_for(human_notes or profile.get("human_notes"), selected))

    # V9 ("flag, never suppress") is recomputed here directly from profile
    # rows rather than trusted from whatever assemble_profile attached to the
    # selected blocks: a negative-equity disclosure must stay visible even if
    # its block never gets a slot, so it cannot depend on slot assignment at
    # all. This mirrors validate_v11/validate_v12's existing render-time
    # recomputation pattern above, and is deliberately a separate line from
    # Note(s)/Source(s) -- see assemble_profile.flagged_items's docstring for
    # why V9 must not be folded into the V3-V6 footnote mechanism.
    advisory_items = flagged_items({"rows": profile.get("rows", ())}, profile.get("segments"))
    advisory_lines = _unique(str(item["message"]) for item in advisory_items)
    advisory_audit_entries = [
        {"slot": "advisory", "block_id": "adv.flagged_items", "std_id": entry.get("std_id"), "doc": entry["doc"], "page": entry.get("page")}
        for item in advisory_items
        for entry in item["provenance"]
    ]

    _set_footer(slide, _unique(rendered_footnotes), _unique(sources), advisory_lines)
    _set_notes(slide, audit_entries + advisory_audit_entries)
    presentation.save(output)

    companion = output.with_suffix(".json")
    companion.write_text(
        json.dumps(
            {
                "profile": profile,
                "render": {
                    "template": str(resolved_template),
                    "slide_number": SLIDE_INDEX + 1,
                    "layout": "1_Blank with Title",
                    "assignments": selected,
                    "figures": audit_entries,
                    "footnotes": _unique(rendered_footnotes),
                    "sources": _unique(sources),
                    "advisory": advisory_lines,
                },
            },
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    return output


def _resolve_template(template: str | Path | None) -> Path:
    candidates = [Path(template)] if template else [
        Path("render") / "templates" / MASTER_FILENAME,
        Path("py") / "render" / "templates" / MASTER_FILENAME,
    ]
    for candidate in candidates:
        if candidate.is_file():
            return candidate.resolve()
    raise RenderError(
        "Ankura master not found. Pass template= or place it at "
        f"render/templates/{MASTER_FILENAME}."
    )


def _validate_assignments(profile: Mapping[str, Any], assignments: Mapping[str, str]) -> None:
    if set(assignments) != set(SLOTS):
        raise RenderError("Exactly one explicit assignment is required for each of the four slots.")
    if len(set(assignments.values())) != len(assignments):
        raise RenderError("One content block cannot be rendered in more than one slot.")
    one_off_evidence = profile.get("lagebericht")
    if one_off_evidence is None and profile.get("one_offs") is not None:
        one_off_evidence = {"one_offs": profile["one_offs"]}
    v11_flags = validate_v11(profile.get("blocks", ()), assignments, lagebericht=one_off_evidence)
    if v11_flags:
        raise RenderError(f"Profile has unresolved V11 flags: {[flag.model_dump() for flag in v11_flags]!r}")
    block_ids = {str(block.get("id")) for block in profile.get("blocks", [])}
    for slot, block_id in assignments.items():
        if block_id not in block_ids:
            raise RenderError(f"{slot} references missing block {block_id!r}.")
        block = next(block for block in profile["blocks"] if block["id"] == block_id)
        if slot not in block.get("eligible_slots", []):
            raise RenderError(f"Block {block_id!r} is not eligible for {slot}.")
        if block.get("unavailable_reason"):
            raise RenderError(f"Block {block_id!r} is unavailable: {block['unavailable_reason']}")
        blocking = [flag for flag in block.get("flags", []) if flag.get("severity") == "blocking"]
        if blocking:
            raise RenderError(f"Block {block_id!r} has unresolved blocking flags: {blocking!r}")
        if "presentation_basis" not in block:
            raise RenderError(f"Rendered block {block_id!r} has no presentation_basis.")
        if not (block.get("unit") or block.get("units")):
            raise RenderError(f"Rendered block {block_id!r} has no unit.")
        for provenance in block.get("provenance", []):
            if not (provenance.get("std_id") or block.get("std_id")) or not provenance.get("doc"):
                raise RenderError(f"Rendered block {block_id!r} has incomplete provenance.")
            if "page" not in provenance:
                raise RenderError(f"Rendered block {block_id!r} provenance must model page, including null.")
        revenue_label = " ".join(str(block.get(key, "")) for key in ("title", "series_label")).lower()
        if (
            block.get("kind") == "chart.column_line"
            and "revenue" in revenue_label
            and block.get("presentation_basis") != "umsatzerloese"
        ):
            raise RenderError("A figure labelled Revenue must use presentation_basis='umsatzerloese' (V1).")

    # V11 above was the only validate.validator rule with a real caller before
    # this; V1-V10 and V12 otherwise only ever ran from tests. blocks double
    # as charted_series here -- each already carries an id, and this is the
    # first point in the render path where slot_assignments (this function's
    # own ``assignments``) are actually known.
    # series_label only -- title is set on every block (including bullet/gap
    # placeholders whose title can itself contain "revenue", e.g. "Revenue
    # Split by Geography"), so falling back to it here would make V1 flag
    # unrelated blocks that were never presented as a Revenue chart.
    blocks_by_id = {str(block.get("id")): block for block in profile.get("blocks", ())}
    axis_labels = {
        slot: blocks_by_id[block_id].get("series_label")
        for slot, block_id in assignments.items() if block_id in blocks_by_id
    }
    result = validate_normalised(
        {"rows": profile.get("rows", ())}, profile.get("segments"),
        charted_series=profile.get("blocks", ()), slot_assignments=assignments,
        axis_labels=axis_labels, lagebericht=one_off_evidence,
    )
    blocking_flags = [flag for flag in result.flags if flag.severity == "blocking"]
    if blocking_flags:
        raise RenderError(f"Profile has unresolved validation flags: {[flag.model_dump() for flag in blocking_flags]!r}")


def _clear_content_regions(slide: Any) -> None:
    """Remove only old content in the four boxes; preserve the OLE shape."""
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            continue
        left, top = shape.left / 914400, shape.top / 914400
        right, bottom = (shape.left + shape.width) / 914400, (shape.top + shape.height) / 914400
        for slot in SLOTS.values():
            # Titles live at y=slot.y. Content starts below the title strip.
            if right >= slot.x and left <= slot.x + slot.width and bottom >= slot.y + 0.32 and top <= slot.y + slot.height:
                shape._element.getparent().remove(shape._element)
                break


def _set_heading(slide: Any, entity: Mapping[str, Any], blocks: Mapping[str, Mapping[str, Any]], assignments: Mapping[str, str]) -> None:
    legal_name = str(entity.get("legal_name", "Company profile"))
    fiscal_year_end = entity.get("fiscal_year_end")
    slide.shapes[1].text = f"At a glance: {legal_name}" + (f" | FYE {fiscal_year_end}" if fiscal_year_end else "")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for slot in SLOTS.values():
            if abs(shape.left / 914400 - slot.x) < 0.08 and abs(shape.top / 914400 - slot.y) < 0.05:
                shape.text = str(blocks[assignments[slot.name]].get("title", ""))


def _render_block(slide: Any, slot: Slot, block: Mapping[str, Any]) -> None:
    kind = block.get("kind")
    if kind in {"chart.column_line", "chart.stacked_column"}:
        _render_native_column_chart(slide, slot, block)
    elif kind in {"bullets", "timeline", "table", "map", "image_grid"}:
        # The block contract presently carries no image/table payload. Render
        # its supplied text only; do not manufacture a proxy visual.
        _render_text_block(slide, slot, block)
    else:
        raise RenderError(f"Unsupported block kind {kind!r} for {block.get('id')!r}.")


def _render_text_block(slide: Any, slot: Slot, block: Mapping[str, Any]) -> None:
    lines = block.get("content") or block.get("bullets") or block.get("summary")
    if not lines:
        reason = block.get("unavailable_reason") or "Content not supplied in profile JSON"
        lines = [reason]
    if isinstance(lines, str):
        lines = [lines]
    box = slide.shapes.add_textbox(Inches(slot.x), Inches(slot.y + 0.43), Inches(slot.width), Inches(slot.height - 0.5))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for i, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.text = str(line)
        paragraph.font.size = Pt(11)
        paragraph.level = 0
        paragraph.text = "• " + paragraph.text


def _render_native_column_chart(slide: Any, slot: Slot, block: Mapping[str, Any]) -> None:
    points = block.get("series")
    if not isinstance(points, Sequence) or not points:
        raise RenderError(f"Chart block {block.get('id')!r} has no series; refusing to draw an empty chart.")
    years, values = _series_values(points, block.get("id", "chart"))
    chart_data = CategoryChartData()
    chart_data.categories = [f"FY'{str(year)[-2:]}A" for year in years]
    chart_data.add_series(str(block.get("series_label", block.get("title", "Value"))), values)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(slot.x + 0.04), Inches(slot.y + 0.50), Inches(slot.width - 0.09), Inches(slot.height - 0.85),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False
    chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.tick_labels.font.size = Pt(7)
    chart.value_axis.tick_labels.number_format = "0"
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = _rgb("003B5C")
    chart.plots[0].series[0].format.line.fill.background()
    chart.plots[0].has_data_labels = False

    # Explicit labels are deliberately generated here, immediately after the
    # chart, from the identical values. They are not chart data labels because
    # the house master uses independently positioned shapes.
    maximum = max(1.0, float(int(max(values) / 50.0 + 1) * 50))
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = maximum
    # The native chart reserves a title/axis margin above its plot area. These
    # coordinates match that area, not the chart frame, so the labels track
    # their bars rather than floating in the title space.
    plot_top = slot.y + 0.68
    plot_height = 1.48
    for index, (year, value) in enumerate(zip(years, values, strict=True)):
        # Native chart plot-area margins put the first category slightly to the
        # right of the frame-based centre calculation.
        x = slot.x + 0.33 + (slot.width - 0.36) * (index + 0.5) / len(values)
        y = plot_top + (maximum - value) / maximum * plot_height - 0.15
        _add_label(slide, _format_millions(value), x - 0.20, max(slot.y + 0.40, y), 0.42, 0.15, 7.5)
        _add_label(slide, f"FY’{str(year)[-2:]}A", x - 0.24, slot.y + slot.height - 0.55, 0.48, 0.15, 7)


def _series_values(points: Sequence[Mapping[str, Any]], block_id: str) -> tuple[list[int], list[float]]:
    years: list[int] = []
    values: list[float] = []
    for point in points:
        try:
            year, value = int(point["fy"]), float(point["value"])
        except (KeyError, TypeError, ValueError) as exc:
            raise RenderError(f"Invalid series point in {block_id!r}: {point!r}") from exc
        years.append(year)
        values.append(value / 1_000_000)
    if years != sorted(years) or len(set(years)) != len(years):
        raise RenderError(f"Series {block_id!r} must have unique fiscal years in ascending order.")
    return years, values


def _add_label(slide: Any, text: str, x: float, y: float, width: float, height: float, font_size: float) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(font_size)
    paragraph.font.name = "Arial"
    paragraph.font.color.rgb = _rgb("000000")


def _set_footer(slide: Any, notes: Sequence[str], sources: Sequence[str], advisory: Sequence[str] = ()) -> None:
    footer = next((shape for shape in slide.shapes if shape.has_text_frame and shape.top / 914400 > 6.9), None)
    if footer is None:
        footer = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.32))
    note_text = " | ".join(notes) if notes else "None"
    source_text = " | ".join(sources) if sources else "Not supplied"
    text = f"Note(s): {note_text}\nSource(s): {source_text}"
    # Deliberately its own line, never merged into Note(s): V9 is advisory and
    # unconditional, not an analyst-authored note against a specific series.
    if advisory:
        text += "\n⚠ Advisory: " + " | ".join(advisory)
    footer.text = text
    for paragraph in footer.text_frame.paragraphs:
        paragraph.font.size = Pt(6.5)


def _set_notes(slide: Any, audit_entries: Sequence[Mapping[str, Any]]) -> None:
    audit_text = "\n".join(
        f"{entry['slot']} | std_id={entry['std_id']} | doc={entry['doc']} | page={entry['page']}"
        for entry in audit_entries
    )
    source_text = "\n".join(
        f"- {entry['doc']}, p. {entry['page'] if entry['page'] is not None else 'n/a'}"
        for entry in _unique_audit_sources(audit_entries)
    )
    slide.notes_slide.notes_text_frame.text = f"[Audit trail]\n{audit_text}\n\n[Sources]\n{source_text}"


def _audit_entries(slot: str, block: Mapping[str, Any]) -> list[dict[str, Any]]:
    entries = []
    for provenance in block.get("provenance", []):
        std_id = provenance.get("std_id") or block.get("std_id")
        if not std_id:
            raise RenderError(f"Rendered block {block.get('id')!r} has provenance without std_id.")
        if not provenance.get("doc"):
            raise RenderError(f"Rendered block {block.get('id')!r} has provenance without doc.")
        entries.append({
            "slot": slot,
            "block_id": block.get("id"),
            "std_id": std_id,
            "doc": provenance["doc"],
            "page": provenance.get("page"),
        })
    if not entries:
        raise RenderError(f"Rendered block {block.get('id')!r} has no provenance.")
    return entries


def _source_lines(block: Mapping[str, Any]) -> list[str]:
    return [
        f"{item['doc']} (p. {item['page'] if item.get('page') is not None else 'n/a'})"
        for item in block.get("provenance", []) if item.get("doc")
    ]


def _flag_notes(block: Mapping[str, Any]) -> list[str]:
    return [str(flag["note"]) for flag in block.get("flags", []) if flag.get("note")]


def _human_notes_for(notes: Any, assignments: Mapping[str, str]) -> list[str]:
    if not notes:
        return []
    if isinstance(notes, str):
        return [notes]
    if isinstance(notes, Mapping):
        selected_ids = set(assignments.values()) | set(assignments)
        result: list[str] = []
        for key, value in notes.items():
            if key in selected_ids:
                result.extend([value] if isinstance(value, str) else value)
        return [str(item) for item in result]
    return [str(item) for item in notes]


def _format_millions(value: float) -> str:
    return f"({abs(value):.1f})" if value < 0 else f"{value:.1f}"


def _unique(items: Iterable[str]) -> list[str]:
    return list(dict.fromkeys(item.strip() for item in items if item and item.strip()))


def _unique_audit_sources(entries: Sequence[Mapping[str, Any]]) -> list[Mapping[str, Any]]:
    unique: dict[tuple[Any, Any], Mapping[str, Any]] = {}
    for entry in entries:
        unique[(entry["doc"], entry["page"])] = entry
    return list(unique.values())


def _rgb(value: str) -> Any:
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(value)
